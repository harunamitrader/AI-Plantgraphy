from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
OUT = ROOT / "docs" / "assets" / "AI-Plantgraphy_illustrated_guide.pptx"

W = Inches(13.333)
H = Inches(7.5)

COLORS = {
    "ink": RGBColor(13, 42, 34),
    "muted": RGBColor(89, 112, 104),
    "green": RGBColor(35, 125, 86),
    "green2": RGBColor(114, 184, 80),
    "mint": RGBColor(232, 244, 236),
    "line": RGBColor(195, 214, 203),
    "cream": RGBColor(252, 251, 244),
    "blue": RGBColor(27, 111, 163),
    "yellow": RGBColor(246, 190, 85),
    "white": RGBColor(255, 255, 255),
}

FONT = "Yu Gothic"


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    if align:
        p.alignment = align
    return box


def add_label(slide, text, x, y, w=1.35, color="green"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.34),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.color.rgb = COLORS[color]
    shape.text_frame.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.name = FONT
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_card(slide, x, y, w, h, fill="white"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    return shape


def add_bullet_lines(slide, lines, x, y, w, h, size=15, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_header(slide, title, kicker=None):
    if kicker:
        add_text(slide, kicker, 0.65, 0.35, 3.2, 0.25, 10, True, "green")
    add_text(slide, title, 0.65, 0.62, 8.4, 0.45, 22, True, "ink")
    slide.shapes.add_picture(str(ASSETS / "ai-plantgraphy-icon.png"), Inches(12.25), Inches(0.28), width=Inches(0.55))


def add_arrow(slide, x1, y1, x2, y2, color="green"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]
    slide.shapes.add_picture(str(ASSETS / "ai-plantgraphy-overview.jpg"), 0, 0, width=W, height=H)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 18
    overlay.line.fill.background()
    add_text(slide, "AI Plantgraphy", 0.8, 0.72, 5.2, 0.55, 30, True, "ink")
    add_text(slide, "植物写真から、自分だけの図鑑を育てるWebアプリ", 0.82, 1.28, 7.2, 0.38, 17, True, "green")
    add_bullet_lines(
        slide,
        [
            "スマホで撮る",
            "自宅PCのGemini CLIで解析する",
            "写真・名前・観察履歴を植物ごとに整理する",
        ],
        0.9,
        5.45,
        7.7,
        1.15,
        16,
    )


def slide_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["mint"]
    add_header(slide, "何ができるツールか", "OVERVIEW")
    items = [
        ("撮影", "庭木・草花を1〜3枚で記録\n連続カメラ、通常カメラ、写真選択に対応"),
        ("解析", "Gemini CLIで種類を推定\n候補と信頼度、不確実な点も保存"),
        ("図鑑", "同じ植物を自動でまとめる\n写真・特徴・手入れメモ・履歴を見返せる"),
    ]
    xs = [0.7, 4.55, 8.4]
    for i, (title, body) in enumerate(items):
        add_card(slide, xs[i], 1.55, 3.35, 4.7)
        add_label(slide, title, xs[i] + 0.25, 1.85, 1.05, ["green", "blue", "yellow"][i])
        add_text(slide, body, xs[i] + 0.28, 2.55, 2.75, 1.6, 16, True if i == 0 else False)
        if i == 0:
            slide.shapes.add_picture(str(ASSETS / "screenshot" / "screenshot8.png"), Inches(xs[i] + 0.28), Inches(4.1), width=Inches(2.7))
        elif i == 1:
            slide.shapes.add_picture(str(ASSETS / "screenshot" / "screenshot4.png"), Inches(xs[i] + 0.28), Inches(4.1), width=Inches(2.7))
        else:
            slide.shapes.add_picture(str(ASSETS / "screenshot" / "screenshot2.png"), Inches(xs[i] + 0.28), Inches(4.1), width=Inches(2.7))
    add_arrow(slide, 4.1, 3.8, 4.45, 3.8)
    add_arrow(slide, 7.95, 3.8, 8.3, 3.8)


def slide_daily_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]
    add_header(slide, "毎日の使い方", "USER FLOW")
    steps = [
        ("1", "追加ページを開く", "スマホのPWAまたはブラウザから開始"),
        ("2", "写真候補を作る", "連続カメラ、通常カメラ、写真から選ぶ"),
        ("3", "1〜3枚を選ぶ", "候補が多いときは後から追加した3枚を初期選択"),
        ("4", "送信して待つ", "画像送信中、種類特定中、解説作成中を表示"),
        ("5", "図鑑で見返す", "植物ごとに写真と観察履歴を整理"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.75 + i * 2.45
        add_card(slide, x, 1.55, 2.08, 3.6, "white")
        add_label(slide, num, x + 0.2, 1.85, 0.5)
        add_text(slide, title, x + 0.22, 2.38, 1.6, 0.35, 15, True)
        add_text(slide, body, x + 0.22, 2.95, 1.62, 1.15, 11, False, "muted")
        if i < 4:
            add_arrow(slide, x + 2.12, 3.3, x + 2.32, 3.3)
    add_text(slide, "基本は3枚。近景・葉・全体などを混ぜると同定しやすくなります。", 0.82, 6.0, 10.2, 0.35, 14, True, "green")


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["mint"]
    add_header(slide, "裏側の仕組み", "SYSTEM")
    nodes = [
        (0.75, 1.65, "スマホ\nWeb / PWA", "撮影・選択・閲覧"),
        (3.15, 1.65, "Tailscale", "外出先から自宅PCへ\nプライベート接続"),
        (5.55, 1.65, "Windows PC\nFastAPI", "受信・軽量化・保存\n画面配信"),
        (8.2, 1.0, "Gemini CLI", "植物の種類を同定\n新規植物の解説生成"),
        (8.2, 3.25, "SQLite + 画像", "観察記録・植物図鑑\n写真ファイルを保存"),
        (10.95, 2.15, "Web図鑑", "スマホで後から閲覧"),
    ]
    for x, y, title, body in nodes:
        add_card(slide, x, y, 1.95, 1.22, "white")
        add_text(slide, title, x + 0.14, y + 0.18, 1.65, 0.42, 13, True)
        add_text(slide, body, x + 0.14, y + 0.66, 1.62, 0.45, 9, False, "muted")
    add_arrow(slide, 2.7, 2.25, 3.08, 2.25)
    add_arrow(slide, 5.1, 2.25, 5.48, 2.25)
    add_arrow(slide, 7.5, 2.05, 8.1, 1.65)
    add_arrow(slide, 7.5, 2.45, 8.1, 3.85)
    add_arrow(slide, 10.15, 3.85, 10.9, 2.85)
    add_text(slide, "外部公開サーバーではなく、自宅PCを中心にしたローカルファースト構成。写真とDBはPC側の data フォルダに残ります。", 0.88, 5.85, 11.4, 0.55, 14, True, "ink")


def slide_data_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]
    add_header(slide, "観察記録と図鑑の分け方", "DATA")
    add_card(slide, 0.8, 1.45, 5.25, 4.55, "white")
    add_card(slide, 7.05, 1.45, 5.25, 4.55, "white")
    add_label(slide, "観察記録", 1.12, 1.78, 1.2, "blue")
    add_text(slide, "その日に撮った写真セット", 1.12, 2.35, 3.6, 0.35, 17, True)
    add_bullet_lines(
        slide,
        ["写真1〜3枚", "撮影日・場所ラベル・メモ", "同定候補と信頼度", "不確実な点、再解析結果"],
        1.15,
        3.0,
        4.2,
        1.6,
        14,
    )
    add_label(slide, "図鑑", 7.38, 1.78, 0.85, "green")
    add_text(slide, "同じ植物をまとめたページ", 7.38, 2.35, 3.6, 0.35, 17, True)
    add_bullet_lines(
        slide,
        ["植物名・学名", "新しい写真から最大12枚", "基本的な特徴", "見た目の特徴と魅力", "手入れメモ、観察履歴"],
        7.42,
        3.0,
        4.2,
        1.8,
        14,
    )
    add_arrow(slide, 6.15, 3.58, 6.95, 3.58)
    add_text(slide, "同じ植物名なら自動で統合", 5.35, 4.0, 2.6, 0.35, 12, True, "green", PP_ALIGN.CENTER)


def slide_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["mint"]
    add_header(slide, "初回セットアップの全体像", "SETUP")
    rows = [
        ("PC", "リポジトリ取得\ninstall_windows.ps1 実行\nGemini CLIをログイン済みにする"),
        ("接続", "PCとスマホでTailscaleにログイン\n設定ページのQRコードをスマホで開く"),
        ("スマホ", "ホーム画面に追加\n設定ページでアプリパスワードとモデルを確認"),
        ("運用", "追加ページから撮影\n図鑑・観察記録・確認待ちを見返す"),
    ]
    for i, (title, body) in enumerate(rows):
        y = 1.45 + i * 1.22
        add_card(slide, 1.0, y, 11.2, 0.85, "white")
        add_label(slide, title, 1.25, y + 0.24, 1.0, ["green", "blue", "yellow", "green"][i])
        add_text(slide, body, 2.55, y + 0.17, 8.8, 0.48, 13, False)
    add_text(slide, "ポイント: ユーザーが入力する「アプリパスワード」はWebアプリの簡易保護用。Gemini CLI側のAPIキーやログインとは別物です。", 1.0, 6.45, 11.2, 0.34, 12, True, "ink")


def slide_oss(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]
    add_header(slide, "OSSとして伝えたいこと", "PUBLIC")
    add_card(slide, 0.8, 1.42, 3.75, 4.8, "white")
    add_card(slide, 4.85, 1.42, 3.75, 4.8, "white")
    add_card(slide, 8.9, 1.42, 3.45, 4.8, "white")
    add_label(slide, "安心", 1.1, 1.75, 0.85, "green")
    add_text(slide, "写真とDBは自宅PCに保存", 1.1, 2.32, 3.0, 0.35, 16, True)
    add_text(slide, "クラウドに植物写真を集約しない設計。外出先接続はTailscaleで閉じた経路を使います。", 1.1, 2.95, 2.95, 1.1, 13, False, "muted")
    add_label(slide, "わかりやすさ", 5.15, 1.75, 1.3, "blue")
    add_text(slide, "非エンジニア向け導線", 5.15, 2.32, 3.0, 0.35, 16, True)
    add_text(slide, "デスクトップショートカット、設定ページ、QRコード、診断画面で迷う場所を減らします。", 5.15, 2.95, 2.95, 1.1, 13, False, "muted")
    add_label(slide, "拡張", 9.2, 1.75, 0.85, "yellow")
    add_text(slide, "Web/PWAから育てる", 9.2, 2.32, 2.5, 0.35, 16, True)
    add_text(slide, "同じAPIを使えば、将来Androidネイティブアプリや別UIも追加できます。", 9.2, 2.95, 2.6, 1.1, 13, False, "muted")
    slide.shapes.add_picture(str(ASSETS / "ai-plantgraphy-header.jpg"), Inches(1.2), Inches(5.0), width=Inches(10.9), height=Inches(1.25))


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_cover(prs)
    slide_value(prs)
    slide_daily_flow(prs)
    slide_architecture(prs)
    slide_data_model(prs)
    slide_setup(prs)
    slide_oss(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
